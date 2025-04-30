import os
from typing import List, Dict, Any
from pptx import Presentation
from .base_processor import DocumentProcessor

class PptxProcessor(DocumentProcessor):
    """
    Document processor for PowerPoint (.pptx) files using python-pptx.
    Extracts text content slide by slide.
    """

    def supports(self, file_path: str) -> bool:
        """
        Checks if the processor supports the given file type.
        """
        return os.path.splitext(file_path)[1].lower() == ".pptx"

    def process(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Processes a .pptx file and extracts content slide by slide.

        Args:
            file_path: The path to the PowerPoint file.

        Returns:
            A list of dictionaries, where each dictionary represents a slide
            with its text content and associated metadata.
        """
        if not self.supports(file_path):
            raise ValueError("File type not supported by PptxProcessor")

        chunks: List[Dict[str, Any]] = []
        try:
            presentation = Presentation(file_path)
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                    # Handle tables in slides if necessary (more complex)
                    # if shape.has_table:
                    #     # Extract table data...

                text_content = "\n".join(slide_text).strip()

                if text_content: # Only process non-empty slides
                    chunks.append({
                        "content": text_content,
                        "type": "slide",
                        "metadata": {
                            "source": file_path,
                            "slide_number": slide_num + 1, # 1-based index
                            # Add other relevant slide properties if needed
                        }
                    })
            return chunks

        except Exception as e:
            print(f"Error processing {file_path}: {e}")
            return []