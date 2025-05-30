import pytest
import os
from docx import Document as DocxDocument # Avoid conflict with test function name
import fitz # PyMuPDF
from pptx import Presentation as PptxPresentation # Avoid conflict

# Assuming processors are in src/processors
from src.processors.base_processor import DocumentProcessor
from src.processors.docx_processor import DocxProcessor
from src.processors.pdf_processor import PdfProcessor
from src.processors.pptx_processor import PptxProcessor

# Define paths for dummy test files
TEST_DIR = "tests/temp_test_files"
DUMMY_DOCX_PATH = os.path.join(TEST_DIR, "dummy.docx")
DUMMY_PDF_PATH = os.path.join(TEST_DIR, "dummy.pdf")
DUMMY_PPTX_PATH = os.path.join(TEST_DIR, "dummy.pptx")
DUMMY_TXT_PATH = os.path.join(TEST_DIR, "dummy.txt")

@pytest.fixture(scope="module", autouse=True)
def setup_and_teardown_test_files():
    """
    Fixture to create and remove dummy test files for processors.
    """
    os.makedirs(TEST_DIR, exist_ok=True)

    # Create dummy docx
    try:
        doc = DocxDocument()
        paragraph = doc.add_paragraph("This is a test paragraph in a docx file with ")
        run = paragraph.add_run("{{ bold }}")
        run.bold = True
        paragraph.add_run(" bold text and ")
        run = paragraph.add_run("{{ italic }}")
        run.italic = True
        paragraph.add_run(" italic text.")
        doc.save(DUMMY_DOCX_PATH)
    except ImportError:
        pytest.skip("python-docx not installed, skipping docx tests")
    except Exception as e:
        print(f"Error creating dummy docx: {e}")

    # Create dummy pdf
    try:
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((10, 50), "This is a test page in a pdf file.")
        doc.save(DUMMY_PDF_PATH)
        doc.close()
    except ImportError:
        pytest.skip("PyMuPDF not installed, skipping pdf tests")
    except Exception as e:
        print(f"Error creating dummy pdf: {e}")

    # Create dummy pptx
    try:
        prs = PptxPresentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Test Presentation"
        prs.save(DUMMY_PPTX_PATH)
    except ImportError:
        pytest.skip("python-pptx not installed, skipping pptx tests")
    except Exception as e:
        print(f"Error creating dummy pptx: {e}")

    # Create dummy txt
    with open(DUMMY_TXT_PATH, "w") as f:
        f.write("This is a plain text file.")

    yield # This is where the tests run

    # Teardown: remove test files and directory
    if os.path.exists(TEST_DIR):
        import shutil
        shutil.rmtree(TEST_DIR)


# --- Test DocxProcessor ---
def test_docx_processor_supports():
    processor = DocxProcessor()
    assert processor.supports(DUMMY_DOCX_PATH) is True
    assert processor.supports(DUMMY_PDF_PATH) is False
    assert processor.supports(DUMMY_TXT_PATH) is False

def test_docx_processor_process():
    processor = DocxProcessor()
    if not os.path.exists(DUMMY_DOCX_PATH):
         pytest.skip("Dummy docx file not created.")
    chunks = processor.process(DUMMY_DOCX_PATH)
    assert isinstance(chunks, list)
    assert len(chunks) > 0
    for chunk in chunks:
        assert isinstance(chunk, dict)
        assert "content" in chunk
        assert "type" in chunk
        assert "metadata" in chunk
        assert isinstance(chunk["content"], str)
        assert isinstance(chunk["metadata"], dict)
        assert "source" in chunk["metadata"]
        assert chunk["metadata"]["source"] == DUMMY_DOCX_PATH

def test_docx_processor_process_formatting():
    processor = DocxProcessor()
    if not os.path.exists(DUMMY_DOCX_PATH):
        pytest.skip("Dummy docx file not created.")
    chunks = processor.process(DUMMY_DOCX_PATH)

    # Find the chunk with bold text
    bold_chunk = next((chunk for chunk in chunks if chunk["content"].strip() == "{{ bold }}"), None)
    assert bold_chunk is not None, "Chunk with '{{ bold }}' not found"
    assert bold_chunk["metadata"]["bold"] is True

    # Find the chunk with italic text
    italic_chunk = next((chunk for chunk in chunks if chunk["content"].strip() == "{{ italic }}"), None)
    assert italic_chunk is not None, "Chunk with '{{ italic }}' not found"
    assert italic_chunk["metadata"]["italic"] is True


def test_docx_processor_process_with_replacements():
    processor = DocxProcessor()
    if not os.path.exists(DUMMY_DOCX_PATH):
        pytest.skip("Dummy docx file not created.")
    replacements = {"bold": "replaced bold text", "italic": "replaced italic text"}
    chunks = processor.process(DUMMY_DOCX_PATH, replacements)

    # Print chunks for debugging
    for chunk in chunks:
        print(chunk)

    # Find the chunk with replaced bold text
    bold_chunk = next((chunk for chunk in chunks if chunk["content"].strip() == "replaced bold text"), None)
    assert bold_chunk is not None
    assert bold_chunk["metadata"]["bold"] is True

    # Find the chunk with replaced italic text
    italic_chunk = next((chunk for chunk in chunks if chunk["content"] == "replaced italic text"), None)
    assert italic_chunk is not None
    assert italic_chunk["metadata"]["italic"] is True


# --- Test PdfProcessor ---
def test_pdf_processor_supports():
    processor = PdfProcessor()
    assert processor.supports(DUMMY_PDF_PATH) is True
    assert processor.supports(DUMMY_DOCX_PATH) is False
    assert processor.supports(DUMMY_TXT_PATH) is False

def test_pdf_processor_process():
    processor = PdfProcessor()
    if not os.path.exists(DUMMY_PDF_PATH):
         pytest.skip("Dummy pdf file not created.")
    chunks = processor.process(DUMMY_PDF_PATH)
    assert isinstance(chunks, list)
    assert len(chunks) > 0
    for chunk in chunks:
        assert isinstance(chunk, dict)
        assert "content" in chunk
        assert "type" in chunk
        assert "metadata" in chunk
        assert isinstance(chunk["content"], str)
        assert isinstance(chunk["metadata"], dict)
        assert chunk["metadata"].get("source") == DUMMY_PDF_PATH
        assert chunk["type"].lower() in ["page", "header"]


# --- Test PptxProcessor ---
def test_pptx_processor_supports():
    processor = PptxProcessor()
    assert processor.supports(DUMMY_PPTX_PATH) is True
    assert processor.supports(DUMMY_DOCX_PATH) is False
    assert processor.supports(DUMMY_TXT_PATH) is False

def test_pptx_processor_process():
    processor = PptxProcessor()
    if not os.path.exists(DUMMY_PPTX_PATH):
         pytest.skip("Dummy pptx file not created.")
    chunks = processor.process(DUMMY_PPTX_PATH)
    assert isinstance(chunks, list)
    assert len(chunks) > 0
    for chunk in chunks:
        assert isinstance(chunk, dict)
        assert "content" in chunk
        assert "type" in chunk
        assert "metadata" in chunk
        assert isinstance(chunk["content"], str)
        assert isinstance(chunk["metadata"], dict)
        assert chunk["metadata"].get("source") == DUMMY_PPTX_PATH
        assert chunk["type"] == "slide"