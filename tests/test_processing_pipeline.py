import pytest
import os
import shutil
from docx import Document as DocxDocument
import fitz # PyMuPDF
from pptx import Presentation as PptxPresentation

# Assuming modules are in src/processors
from src.processors.document_parser import parse_document
from src.processors.chunking import chunk_document_elements

# Define paths for dummy test files
TEST_DIR = "tests/temp_processing_test_files"
DUMMY_DOCX_PATH = os.path.join(TEST_DIR, "dummy_parse_chunk.docx")
DUMMY_PDF_PATH = os.path.join(TEST_DIR, "dummy_parse_chunk.pdf")
DUMMY_PPTX_PATH = os.path.join(TEST_DIR, "dummy_parse_chunk.pptx")
DUMMY_TXT_PATH = os.path.join(TEST_DIR, "dummy_parse_chunk.txt")
DUMMY_EMPTY_TXT_PATH = os.path.join(TEST_DIR, "dummy_empty.txt")

@pytest.fixture(scope="module", autouse=True)
def setup_and_teardown_test_files():
    """
    Fixture to create and remove dummy test files for processing pipeline tests.
    """
    os.makedirs(TEST_DIR, exist_ok=True)

    # Create dummy docx
    try:
        doc = DocxDocument()
        doc.add_paragraph("This is the first paragraph for parsing and chunking.")
        doc.add_paragraph("This is the second paragraph, which is a bit longer to test chunking boundaries.")
        doc.add_paragraph("Third paragraph.")
        doc.save(DUMMY_DOCX_PATH)
    except ImportError:
        print("python-docx not installed, skipping dummy docx creation.")
    except Exception as e:
        print(f"Error creating dummy docx: {e}")

    # Create dummy pdf
    try:
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((10, 50), "This is the first page for parsing and chunking.")
        page.insert_text((10, 100), "This is the second sentence on the first page.")
        page2 = doc.new_page()
        page2.insert_text((10, 50), "This is the second page.")
        doc.save(DUMMY_PDF_PATH)
        doc.close()
    except ImportError:
        print("PyMuPDF not installed, skipping dummy pdf creation.")
    except Exception as e:
        print(f"Error creating dummy pdf: {e}")

    # Create dummy pptx
    try:
        prs = PptxPresentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Test Presentation for Parsing"
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "This is the body text on the first slide."
        prs.save(DUMMY_PPTX_PATH)
    except ImportError:
        print("python-pptx not installed, skipping dummy pptx creation.")
    except Exception as e:
        print(f"Error creating dummy pptx: {e}")

    # Create dummy txt
    with open(DUMMY_TXT_PATH, "w") as f:
        f.write("This is a plain text file for parsing and chunking.\n")
        f.write("It has multiple lines.\n")
        f.write("And a third line.")

    # Create dummy empty txt
    with open(DUMMY_EMPTY_TXT_PATH, "w") as f:
        f.write("")


    yield # This is where the tests run

    # Teardown: remove test files and directory
    if os.path.exists(TEST_DIR):
        shutil.rmtree(TEST_DIR)


# --- Test parse_document function ---
def test_parse_document_docx():
    if not os.path.exists(DUMMY_DOCX_PATH):
         pytest.skip("Dummy docx file not created.")
    elements = parse_document(DUMMY_DOCX_PATH)
    assert isinstance(elements, list)
    assert len(elements) > 0
    for element in elements:
        assert isinstance(element, dict)
        assert "content" in element
        assert "type" in element
        assert "metadata" in element
        assert isinstance(element["content"], str)
        assert isinstance(element["metadata"], dict)
        assert element["metadata"].get("source") == DUMMY_DOCX_PATH

def test_parse_document_pdf():
    if not os.path.exists(DUMMY_PDF_PATH):
         pytest.skip("Dummy pdf file not created.")
    elements = parse_document(DUMMY_PDF_PATH)
    assert isinstance(elements, list)
    assert len(elements) > 0
    for element in elements:
        assert isinstance(element, dict)
        assert "content" in element
        assert "type" in element
        assert "metadata" in element
        assert isinstance(element["content"], str)
        assert isinstance(element["metadata"], dict)
        assert element["metadata"].get("source") == DUMMY_PDF_PATH
        assert element["type"] == "page" # unstructured often categorizes PDF pages as 'page'

def test_parse_document_pptx():
    if not os.path.exists(DUMMY_PPTX_PATH):
         pytest.skip("Dummy pptx file not created.")
    elements = parse_document(DUMMY_PPTX_PATH)
    assert isinstance(elements, list)
    assert len(elements) > 0
    for element in elements:
        assert isinstance(element, dict)
        assert "content" in element
        assert "type" in element
        assert "metadata" in element
        assert isinstance(element["content"], str)
        assert isinstance(element["metadata"], dict)
        assert element["metadata"].get("source") == DUMMY_PPTX_PATH
        # unstructured might categorize pptx elements differently, check for common types
        assert element["type"] in ["title", "narrative_text", "text"]


def test_parse_document_txt():
    elements = parse_document(DUMMY_TXT_PATH)
    assert isinstance(elements, list)
    assert len(elements) > 0
    for element in elements:
        assert isinstance(element, dict)
        assert "content" in element
        assert "type" in element
        assert "metadata" in element
        assert isinstance(element["content"], str)
        assert isinstance(element["metadata"], dict)
        assert element["metadata"].get("source") == DUMMY_TXT_PATH
        assert element["type"] == "text" # unstructured often categorizes text as 'text'

def test_parse_document_empty():
    elements = parse_document(DUMMY_EMPTY_TXT_PATH)
    assert isinstance(elements, list)
    assert len(elements) == 0 # Should return empty list for empty file

def test_parse_document_non_existent():
    elements = parse_document("non_existent_file.xyz")
    assert isinstance(elements, list)
    assert len(elements) == 0 # Should return empty list for non-existent file


# --- Test chunk_document_elements function ---
def test_chunk_document_elements_basic():
    elements = [
        {"content": "This is the first sentence.", "type": "text", "metadata": {"page": 1}},
        {"content": "This is the second sentence.", "type": "text", "metadata": {"page": 1}},
        {"content": "This is the third sentence.", "type": "text", "metadata": {"page": 2}},
    ]
    # Test with chunk size smaller than content
    chunks = chunk_document_elements(elements, chunk_size=20, chunk_overlap=5)
    assert isinstance(chunks, list)
    assert len(chunks) > 0
    for chunk in chunks:
        assert isinstance(chunk, dict)
        assert "content" in chunk
        assert "metadata" in chunk
        assert isinstance(chunk["content"], str)
        assert isinstance(chunk["metadata"], dict)
        assert len(chunk["content"]) <= 20 + len("\n") # Allow for added newline

def test_chunk_document_elements_overlap():
    elements = [
        {"content": "This is a long sentence that will be split into multiple chunks.", "type": "text", "metadata": {"page": 1}},
    ]
    chunk_size = 30
    chunk_overlap = 10
    chunks = chunk_document_elements(elements, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    assert isinstance(chunks, list)
    assert len(chunks) > 1 # Should be split
    # Check overlap
    assert chunks[0]["content"][-chunk_overlap:] in chunks[1]["content"]

def test_chunk_document_elements_empty_input():
    chunks = chunk_document_elements([])
    assert isinstance(chunks, list)
    assert len(chunks) == 0

def test_chunk_document_elements_large_element_single_chunk():
    elements = [
        {"content": "A" * 100, "type": "text", "metadata": {"page": 1}},
    ]
    chunks = chunk_document_elements(elements, chunk_size=200, chunk_overlap=0)
    assert len(chunks) == 1
    assert chunks[0]["content"] == "A" * 100

def test_chunk_document_elements_large_element_multiple_chunks():
    elements = [
        {"content": "A" * 100, "type": "text", "metadata": {"page": 1}},
    ]
    chunks = chunk_document_elements(elements, chunk_size=30, chunk_overlap=5)
    assert len(chunks) > 1
    for chunk in chunks:
        assert len(chunk["content"]) <= 30 + len("\n") # Allow for added newline

def test_chunk_document_elements_metadata_inheritance():
    elements = [
        {"content": "Chunk 1 content.", "type": "text", "metadata": {"page": 1, "source": "doc.txt"}},
        {"content": "Chunk 2 content.", "type": "text", "metadata": {"page": 1, "source": "doc.txt"}},
    ]
    chunks = chunk_document_elements(elements, chunk_size=50, chunk_overlap=0)
    assert len(chunks) == 2
    assert chunks[0]["metadata"] == {"page": 1, "source": "doc.txt"}
    assert chunks[1]["metadata"] == {"page": 1, "source": "doc.txt"} # Metadata should be inherited/updated